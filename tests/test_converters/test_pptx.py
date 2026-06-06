import tempfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from core.converters.pptx import convert_pptx

def _make_pptx(tmp_dir: Path) -> Path:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide One"
    slide.placeholders[1].text = "Bullet point one"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "Speaker note here"
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Slide Two"
    path = tmp_dir / "sample.pptx"
    prs.save(str(path))
    return path

def test_pptx_slide_titles():
    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        path = _make_pptx(tmp)
        md, assets = convert_pptx(path, tmp / "assets")
        assert "# Slide One" in md
        assert "# Slide Two" in md

def test_pptx_slide_separator():
    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        path = _make_pptx(tmp)
        md, assets = convert_pptx(path, tmp / "assets")
        assert "---" in md

def test_pptx_speaker_notes():
    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        path = _make_pptx(tmp)
        md, assets = convert_pptx(path, tmp / "assets")
        assert "> Note: Speaker note here" in md

def test_pptx_returns_asset_list():
    with tempfile.TemporaryDirectory() as tmp:
        tmp = Path(tmp)
        path = _make_pptx(tmp)
        md, assets = convert_pptx(path, tmp / "assets")
        assert isinstance(assets, list)

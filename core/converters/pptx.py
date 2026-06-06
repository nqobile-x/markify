from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from core.extractor import save_image

def _bullets(text_frame) -> list[str]:
    lines = []
    for para in text_frame.paragraphs:
        level = para.level
        text = para.text.strip()
        if not text:
            continue
        prefix = "  " * level + "- "
        lines.append(prefix + text)
    return lines

def convert_pptx(source: Path, assets_dir: Path) -> tuple[str, list[Path]]:
    prs = Presentation(str(source))
    sections: list[str] = []
    asset_paths: list[Path] = []
    img_index = 1

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []

        if slide_num > 1:
            slide_lines.append("---")

        # Title
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_lines.append(f"# {slide.shapes.title.text.strip()}")

        # Body shapes
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_bytes = shape.image.blob
                    ext = shape.image.ext
                    path = save_image(image_bytes, assets_dir, img_index, ext)
                    asset_paths.append(path)
                    slide_lines.append(f"![image_{img_index}](assets/image_{img_index}.{ext})")
                    img_index += 1
                except Exception:
                    slide_lines.append("<!-- image could not be extracted -->")
            elif shape.has_text_frame:
                slide_lines.extend(_bullets(shape.text_frame))

        # Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            # python-pptx may include a default prompt placeholder in the first paragraph;
            # collect only paragraphs that belong to the actual notes text frame body
            notes_text = notes_tf.text.strip()
            if notes_text:
                slide_lines.append(f"> Note: {notes_text}")

        sections.append("\n".join(slide_lines))

    return "\n\n".join(sections), asset_paths
